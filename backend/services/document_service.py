from io import BytesIO
import subprocess
import tempfile
import os
import asyncio
from concurrent.futures import ThreadPoolExecutor

# All conversions here are CPU-bound (parsing, rendering, subprocess calls)
# and must never block the asyncio event loop, or one slow conversion would
# stall every other concurrent request. Everything routes through this pool.
_executor = ThreadPoolExecutor(max_workers=4)

# Some deployment environments install LibreOffice at a non-default path;
# this lets it be overridden without touching code.
LIBREOFFICE_BIN = os.environ.get("LIBREOFFICE_PATH", "soffice")


async def _run_blocking(fn, *args):
    """Run a CPU-bound/blocking function off the asyncio event loop."""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(_executor, fn, *args)


class DocumentService:
    """Document format conversions (Word, Excel, PowerPoint, PDF).

    PDF -> Office conversions use pdf2docx / pdfplumber / PyMuPDF, which
    reconstruct real layout (text position, tables, images) rather than
    dumping plain text into a document. Office -> PDF conversions shell out
    to LibreOffice headless, which is the only open-source approach that
    reliably preserves original formatting (fonts, styles, page layout,
    headers/footers) for arbitrary Office documents — no Python library
    does this as faithfully as the actual originating application's own
    rendering engine.
    """

    # ------------------------------------------------------------------
    # PDF -> Word
    # ------------------------------------------------------------------
    @staticmethod
    async def pdf_to_word(pdf_content: bytes) -> bytes:
        """Convert PDF to a real, formatted .docx via pdf2docx.

        pdf2docx parses each page's layout (text blocks, tables, images)
        and reconstructs it as native Word content, preserving text
        positioning and tables far better than extracting plain text.
        """
        return await _run_blocking(DocumentService._pdf_to_word_sync, pdf_content)

    @staticmethod
    def _pdf_to_word_sync(pdf_content: bytes) -> bytes:
        from pdf2docx import Converter

        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, "input.pdf")
            output_path = os.path.join(temp_dir, "output.docx")

            with open(input_path, "wb") as f:
                f.write(pdf_content)

            try:
                cv = Converter(input_path)
                try:
                    cv.convert(output_path)
                finally:
                    cv.close()
            except Exception as e:
                raise Exception(f"PDF to Word conversion failed: {str(e)}")

            if not os.path.exists(output_path):
                raise Exception("PDF to Word conversion produced no output")

            with open(output_path, "rb") as f:
                return f.read()

    # ------------------------------------------------------------------
    # PDF -> Excel
    # ------------------------------------------------------------------
    @staticmethod
    async def pdf_to_excel(pdf_content: bytes) -> bytes:
        """Convert PDF tables to a formatted .xlsx via pdfplumber's table
        detection. Pages with no detected table fall back to one row of
        text per line so content is never silently dropped."""
        return await _run_blocking(DocumentService._pdf_to_excel_sync, pdf_content)

    @staticmethod
    def _pdf_to_excel_sync(pdf_content: bytes) -> bytes:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        import pdfplumber

        wb = Workbook()
        wb.remove(wb.active)

        with pdfplumber.open(BytesIO(pdf_content)) as pdf:
            table_count = 0
            for page_num, page in enumerate(pdf.pages, 1):
                tables = page.extract_tables()

                if tables:
                    for table_data in tables:
                        if not table_data:
                            continue
                        table_count += 1
                        sheet_name = f"Table {table_count}"[:31]  # Excel sheet name limit
                        ws = wb.create_sheet(title=sheet_name)

                        for row_idx, row_data in enumerate(table_data, 1):
                            for col_idx, cell_data in enumerate(row_data, 1):
                                cell = ws.cell(row=row_idx, column=col_idx)
                                cell.value = cell_data
                                if row_idx == 1:
                                    cell.font = Font(bold=True)
                                    cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
                                    cell.alignment = Alignment(horizontal="center")

                        for col in ws.columns:
                            max_length = max((len(str(c.value)) for c in col if c.value), default=0)
                            ws.column_dimensions[col[0].column_letter].width = min(max_length + 2, 50)
                else:
                    text = page.extract_text()
                    if text:
                        sheet_name = f"Page {page_num}"[:31]
                        ws = wb.create_sheet(title=sheet_name)
                        for row_idx, line in enumerate(text.split("\n"), 1):
                            ws.cell(row=row_idx, column=1).value = line

        if not wb.sheetnames:
            wb.create_sheet(title="Sheet1")

        output = BytesIO()
        wb.save(output)
        return output.getvalue()

    # ------------------------------------------------------------------
    # PDF -> PowerPoint
    # ------------------------------------------------------------------
    @staticmethod
    async def pdf_to_pptx(pdf_content: bytes) -> bytes:
        """Convert each PDF page into a full-bleed slide image.

        PDFs have no slide/shape semantics, so there's no reliable way to
        reconstruct individually-editable text boxes on a slide without a
        full layout-reconstruction engine. Rendering each page as a
        high-resolution image and placing it full-bleed on its own slide
        preserves the exact visual appearance (fonts, layout, images) —
        the realistic scope for "PDF to PowerPoint" without a much larger
        engine. Slide content will not be individually editable text.
        """
        return await _run_blocking(DocumentService._pdf_to_pptx_sync, pdf_content)

    @staticmethod
    def _pdf_to_pptx_sync(pdf_content: bytes) -> bytes:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Emu

        EMU_PER_INCH = 914400
        RENDER_DPI = 150

        doc = fitz.open(stream=pdf_content, filetype="pdf")
        try:
            if doc.page_count == 0:
                raise Exception("PDF has no pages")

            prs = Presentation()
            zoom = RENDER_DPI / 72
            matrix = fitz.Matrix(zoom, zoom)

            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=matrix)
                img_bytes = pix.tobytes("png")

                slide_width = Emu(int(page.rect.width / 72 * EMU_PER_INCH))
                slide_height = Emu(int(page.rect.height / 72 * EMU_PER_INCH))

                if i == 0:
                    prs.slide_width = slide_width
                    prs.slide_height = slide_height

                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
                slide.shapes.add_picture(BytesIO(img_bytes), 0, 0, width=prs.slide_width, height=prs.slide_height)
        finally:
            doc.close()

        output = BytesIO()
        prs.save(output)
        return output.getvalue()

    # ------------------------------------------------------------------
    # Office -> PDF (all via LibreOffice headless)
    # ------------------------------------------------------------------
    @staticmethod
    async def word_to_pdf(docx_content: bytes) -> bytes:
        """Convert a Word document to PDF via LibreOffice headless — the
        only open-source approach that reliably preserves fonts, styles
        and layout for arbitrary .docx input."""
        return await DocumentService._convert_with_libreoffice(docx_content, "docx", "pdf")

    @staticmethod
    async def excel_to_pdf(xlsx_content: bytes) -> bytes:
        """Convert an Excel workbook to PDF via LibreOffice headless."""
        return await DocumentService._convert_with_libreoffice(xlsx_content, "xlsx", "pdf")

    @staticmethod
    async def pptx_to_pdf(pptx_content: bytes) -> bytes:
        """Convert a PowerPoint deck to PDF via LibreOffice headless."""
        return await DocumentService._convert_with_libreoffice(pptx_content, "pptx", "pdf")

    @staticmethod
    async def _convert_with_libreoffice(content: bytes, input_ext: str, output_format: str) -> bytes:
        return await _run_blocking(
            DocumentService._convert_with_libreoffice_sync, content, input_ext, output_format
        )

    @staticmethod
    def _convert_with_libreoffice_sync(content: bytes, input_ext: str, output_format: str) -> bytes:
        """Convert a document using LibreOffice's headless CLI.

        The input file must keep its real extension (.docx/.xlsx/.pptx) so
        LibreOffice's import filter can detect the format — writing it as a
        generic name makes detection unreliable. LibreOffice also names its
        output after the input's base name (input.docx -> input.pdf), so
        the expected output path must match that exactly.
        """
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, f"input.{input_ext}")
            output_path = os.path.join(temp_dir, f"input.{output_format}")

            with open(input_path, "wb") as f:
                f.write(content)

            cmd = [
                LIBREOFFICE_BIN,
                "--headless",
                "--norestore",
                "--convert-to", output_format,
                "--outdir", temp_dir,
                input_path,
            ]

            try:
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            except FileNotFoundError:
                raise Exception(
                    "LibreOffice is not installed on this server. Office document conversion requires it."
                )
            except subprocess.TimeoutExpired:
                raise Exception("Conversion timed out — the document may be too large or complex.")

            if result.returncode != 0 or not os.path.exists(output_path):
                stderr = (result.stderr or "").strip()
                raise Exception(f"LibreOffice conversion failed{f': {stderr}' if stderr else ''}")

            with open(output_path, "rb") as f:
                return f.read()
